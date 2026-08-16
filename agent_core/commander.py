"""
commander.py — 仙人掌 Agent 主控制器
- 总工程师：调度工具 + 子代理 + 记忆
- 支持 continue/remember/recall 指令
- 自动摘要提醒
"""
import asyncio
import time
import logging
import uuid
import re
import os
import sys
import json
from dataclasses import dataclass
from enum import Enum
from typing import Optional, Callable
from types import SimpleNamespace

from .protocol import Protocol, ExecutionResult
from .session import DeepSeekSession, SessionConfig, MessageRole
# 注意：子代理实际实现在 subagent_manager.py（同浏览器子窗口 + 完整 Commander 循环），
# 由 _browser_research/_browser_visit 工具通过 get_subagent_manager 使用。旧的 subagent.py
# 仅保留 BrowserSubAgent 基类，已不再被本模块引用，故此处不导入它。
from .memory_manager import MemoryManager
# 动态平台注册表（适配更多网页 AI）：从 platforms.json / platforms.user.json 读取
from agent_core.platform_browser import list_platforms as _list_platforms

logger = logging.getLogger("commander")


class EventType(Enum):
    THINKING = "thinking"
    TOOL_START = "tool_start"
    TOOL_END = "tool_end"
    TOOL_ERROR = "tool_error"
    COMMAND_DETECTED = "command_detected"
    COMMAND_EXECUTING = "command_executing"
    COMMAND_SUCCESS = "command_success"
    COMMAND_ERROR = "command_error"
    AI_FINAL_REPLY = "ai_final_reply"
    MEMORY_REMINDER = "memory_reminder"
    CONTINUE_READY = "continue_ready"
    CORRECTION_SENT = "correction_sent"
    ERROR = "error"
    QUESTION = "question"
    SESSION_UPDATED = "session_updated"


@dataclass
class AgentEvent:
    event_type: EventType
    data: any = None


# ============================================================
# 工具注册表
# ============================================================
class ToolRegistry:
    """工具注册表，支持子代理工具"""

    def __init__(self, commander):
        self._commander = commander
        self._tools: dict = {}

    def register(self, name: str, description: str, fn: Callable):
        self._tools[name] = {"description": description, "fn": fn}

    def list_tools(self) -> list:
        return [
            {"name": name, "description": info["description"]}
            for name, info in self._tools.items()
        ]

    # 常见「模型臆想的工具名」→ 真实工具 的别名映射。
    # 浏览器 AI 不一定严格照搬注册名，常凭语义编造 generate_word_report / make_ppt 等，
    # 这里做一层兜底，避免「未知工具」直接把整条任务判死。
    TOOL_ALIASES = {
        # Word / PPT
        "generate_word_report": "docx_create", "make_word": "docx_create",
        "create_word": "docx_create", "word_create": "docx_create",
        "make_docx": "docx_create", "docx": "docx_create",
        "generate_docx": "docx_create", "create_docx": "docx_create",
        "generate_ppt": "pptx_create", "make_ppt": "pptx_create",
        "ppt_create": "pptx_create", "pptx": "pptx_create",
        "create_ppt": "pptx_create",
        # 文件读写
        "write_file": "file_write", "create_file": "file_write",
        "save_file": "file_write", "write_text": "file_write",
        "read_file": "file_read", "read_text": "file_read",
        "edit_file": "file_edit", "modify_file": "file_edit",
        "append_file": "file_edit", "update_file": "file_edit",
        "make_dir": "dir_create", "mkdir": "dir_create",
        "delete_file": "file_delete", "remove_file": "file_delete",
        "rm": "file_delete", "ls": "file_list", "list_dir": "file_list",
        # 搜索
        "search": "grep", "grep_search": "grep", "find": "glob",
        "search_files": "glob", "fetch": "web_fetch", "fetch_url": "web_fetch",
        "search_web": "websearch", "web_search": "websearch",
        # 任务 / 记忆
        "todo": "todowrite", "todo_list": "todoread", "read_todo": "todoread",
        "remember_info": "remember", "recall_info": "recall",
        # 浏览器
        "click": "browser_click", "fill": "browser_fill",
        "screenshot": "browser_screenshot", "search_browser": "browser_search",
        # 子代理
        "research": "browser_research", "visit": "browser_visit",
        # 收尾（模型常用 answer / respond / reply 给最终答复，统一映射到 done）
        "answer": "done", "respond": "done", "reply": "done",
        "final_answer": "done", "finish": "done",
    }

    def _resolve_tool_name(self, raw: str):
        """把模型传来的（可能不对的）工具名解析成真实注册名。
        返回 (resolved_name, mapped_note)。mapped_note 为空表示原名即合法。"""
        import difflib
        name = (raw or "").strip()
        if name in self._tools:
            return name, ""
        # 1) 显式别名
        mapped = self.TOOL_ALIASES.get(name.lower())
        if mapped is None:
            # 2) 模糊匹配（编辑距离）
            matches = difflib.get_close_matches(
                name.lower(), [t.lower() for t in self._tools], n=1, cutoff=0.6
            )
            if matches:
                for t in self._tools:
                    if t.lower() == matches[0]:
                        mapped = t
                        break
        if mapped is not None and mapped in self._tools:
            return mapped, f"工具名 '{raw}' 未注册，已自动映射到 '{mapped}'"
        return name, ""

    async def execute(self, tool_name: str, params: dict) -> "ExecutionResult":
        """异步执行工具（由 Commander 调用）"""
        import uuid
        resolved, note = self._resolve_tool_name(tool_name)
        if resolved not in self._tools:
            return ExecutionResult(id=str(uuid.uuid4()), status="error",
                                  error=f"未知工具: {tool_name}", tool=tool_name)
        fn = self._tools[resolved]["fn"]
        try:
            result = await fn(**params) if asyncio.iscoroutinefunction(fn) else fn(**params)
            out = str(result)
            if note:
                out = note + "\n" + out
            return ExecutionResult(id=str(uuid.uuid4()), status="success", output=out, tool=resolved)
        except Exception as e:
            return ExecutionResult(id=str(uuid.uuid4()), status="error", error=str(e), tool=resolved)


# ============================================================
# 主控制器 Commander
# ============================================================

class Commander:
    """
    Agent 主控制器（总工程师模式）
    - 对话管理：维护 session + 事件回调
    - 工具执行：调用工具注册表
    - 子代理管理：维护 SubAgentManager（任务状态 + 凭据共享）
    - 记忆管理：自动摘要 + 长期记忆（remember/recall）
    - 错误纠正：commander fix
    """

    def __init__(
        self,
        browser_manager,
        session: Optional[DeepSeekSession] = None,
        work_dir: str = "",
        on_event: Optional[Callable[[AgentEvent], None]] = None,
    ):
        self._bm = browser_manager
        self._session = session
        self._work_dir = work_dir or os.getcwd()
        self._on_event = on_event
        self._tools = ToolRegistry(self)
        self._subagent_manager = None  # 实际由 subagent_manager.get_subagent_manager 提供
        self._memory: Optional[MemoryManager] = None
        self._history_turns: int = 0
        self._accumulated_prompt: str = ""
        self._correction_pending: Optional[str] = None
        self._running = False
        self._protocol = Protocol()
        self._pending_attachments: list = []  # 待附带给 AI 平台的文件（图片/PDF等）

        # 先注册工具，再构建系统提示词（提示词依赖工具列表）
        self._register_tools()
        self._system_prompt = self._build_system_prompt()

    def _register_tools(self):
        """注册所有工具（包括子代理工具）"""
        import os
        from pathlib import Path

        # ─── 文件操作工具 ───
        async def file_write(**params):
            p = Path(params.get("path", ""))
            content = params.get("content", "")
            full_path = self._safe_path(p)
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text(content, encoding="utf-8")
            return f"文件已写入: {full_path}"

        async def file_read(**params):
            """读取文件（对标 OpenCode read，支持行范围）。"""
            p = Path(params.get("path", ""))
            full_path = self._safe_path(p)
            if not full_path.exists():
                return f"错误: 文件不存在 {full_path}"
            text = full_path.read_text(encoding="utf-8")
            all_lines = text.split("\n")
            offset = int(params.get("offset", 1))
            limit = params.get("limit")
            if offset < 1:
                offset = 1
            seg = all_lines[offset - 1:]
            if limit is not None and int(limit) > 0:
                seg = seg[:int(limit)]
            numbered = "\n".join(f"{offset + i:>6}\t{ln}" for i, ln in enumerate(seg))
            return f"[文件 {full_path} 行 {offset}-{offset + len(seg) - 1} / 共 {len(all_lines)} 行]\n{numbered}"

        async def file_list(**params):
            p = Path(params.get("path", "."))
            full_path = self._safe_path(p)
            if not full_path.exists():
                return "目录不存在"
            items = list(full_path.iterdir())
            return "\n".join([f"{'[DIR]' if i.is_dir() else '[FILE]'} {i.name}" for i in items])

        async def dir_create(**params):
            raw = str(params.get("path", "")).strip().strip('"').strip("'").strip()
            if not raw:
                return "错误: 未提供 path"
            p = Path(raw)
            full_path = self._safe_path(p)
            full_path.mkdir(parents=True, exist_ok=True)
            return f"目录已创建: {full_path}"

        async def file_delete(**params):
            p = Path(params.get("path", ""))
            full_path = self._safe_path(p)
            if full_path.exists():
                if full_path.is_file():
                    full_path.unlink()
                else:
                    import shutil
                    shutil.rmtree(full_path)
                return f"已删除: {full_path}"
            return f"文件不存在: {full_path}"

        self._tools.register("file_write", "写入文件", file_write)
        self._tools.register("write", "写入文件(OpenCode 规范名, 同 file_write)", file_write)
        self._tools.register("file_read", "读取文件(支持 offset/limit 行范围)", file_read)
        self._tools.register("read", "读取文件(OpenCode 规范名, 同 file_read)", file_read)
        self._tools.register("file_list", "列出目录", file_list)
        self._tools.register("list", "列出目录(OpenCode 规范名, 同 file_list)", file_list)
        self._tools.register("dir_create", "创建目录", dir_create)
        self._tools.register("file_delete", "删除文件/目录", file_delete)

        # ─── 对标 OpenCode 文件编辑工具 ───
        async def file_edit(**params):
            """编辑文件（对标 OpenCode edit，支持多种模式）。
            params:
              path: 文件路径
              mode: 编辑模式，默认 "replace"
                    - "replace"（默认）: 把 old 字符串替换为 new；或 pattern/repl 正则替换
                    - "append": 在文件末尾追加 text（自动补换行）
                    - "insert": 在第 line 行(1-based)前插入 text；after=true 则在 line 行后插入
                    - "delete": 删除 [start,end] 行范围(1-based 闭区间)；也可配合 old 删除匹配行
              old / new: replace 模式的待替换文本
              pattern / repl: 正则替换（优先于 old/new）
              all: replace 正则模式下是否替换全部（默认仅第一处）
              text: append/insert 模式的插入文本
              line / start / end: 行号参数
            """
            import re
            p = Path(params.get("path", ""))
            full_path = self._safe_path(p)
            if not full_path.exists():
                return f"错误: 文件不存在 {full_path}"
            mode = (params.get("mode") or "replace").lower()
            raw = full_path.read_text(encoding="utf-8")
            lines = raw.split("\n")

            if mode == "replace":
                old = params.get("old")
                new = params.get("new", "")
                pattern = params.get("pattern")
                if pattern:
                    flags = re.MULTILINE
                    if params.get("dotall"):
                        flags |= re.DOTALL
                    count = 0 if params.get("all") else 1
                    text2, n = re.subn(pattern, params.get("repl", new), raw, count=count, flags=flags)
                    if n == 0:
                        return f"未找到匹配正则: {pattern}"
                else:
                    if old is None:
                        return "错误: 必须提供 old 或 pattern"
                    if old not in raw:
                        return f"错误: 未找到待替换文本: {old[:80]}"
                    n = raw.count(old)
                    if n > 1 and not params.get("all"):
                        return f"错误: old 出现 {n} 次，存在歧义，请补充上下文或加 all=true 全量替换"
                    text2 = raw.replace(old, new, 0) if params.get("all") else raw.replace(old, new, 1)
                full_path.write_text(text2, encoding="utf-8")
                return f"已编辑 {full_path}（替换 {n} 处，字符数 {len(raw)}→{len(text2)}）"

            elif mode == "append":
                text = params.get("text", "")
                content = (raw + "\n" if not raw.endswith("\n") and raw else raw) + text
                if not content.endswith("\n"):
                    content += "\n"
                full_path.write_text(content, encoding="utf-8")
                return f"已追加到 {full_path}（行数 {len(lines)}→{len(content.split(chr(10)))})"

            elif mode == "insert":
                text = params.get("text", "")
                lineno = int(params.get("line", len(lines) + 1))
                if lineno < 1:
                    lineno = 1
                if params.get("after"):
                    lineno = min(lineno + 1, len(lines) + 1)
                lines.insert(lineno - 1, text)
                full_path.write_text("\n".join(lines), encoding="utf-8")
                return f"已在第 {lineno} 行插入内容: {full_path}（行数 {len(raw.split(chr(10)))}→{len(lines)}）"

            elif mode == "delete":
                old = params.get("old")
                if old is not None:
                    kept = [ln for ln in lines if old not in ln]
                    removed = len(lines) - len(kept)
                    if removed == 0:
                        return f"未找到含待删文本的行: {old[:60]}"
                    full_path.write_text("\n".join(kept), encoding="utf-8")
                    return f"已删除 {removed} 行（含 '{old[:40]}'）: {full_path}"
                start = int(params.get("start", 1))
                end = int(params.get("end", start))
                if start < 1:
                    start = 1
                if end > len(lines):
                    end = len(lines)
                if start > end:
                    return f"错误: start({start}) > end({end})"
                del lines[start - 1:end]
                full_path.write_text("\n".join(lines), encoding="utf-8")
                return f"已删除第 {start}-{end} 行: {full_path}（行数 {len(raw.split(chr(10)))}→{len(lines)}）"

            else:
                return f"错误: 未知 mode={mode}（支持 replace/append/insert/delete）"

        async def grep(**params):
            """在目录中按正则/文本搜索文件内容（对标 opencode/rg 的代码搜索）。
            params:
              pattern: 搜索模式（正则）
              path: 搜索根目录，默认工作目录
              glob: 可选文件名过滤，如 '*.py'
              max_results: 可选，默认 50
              ignore_case: 可选 bool
            """
            import re
            root = self._safe_path(Path(params.get("path", ".")))
            pattern = params.get("pattern", "")
            if not pattern:
                return "错误: 必须提供 pattern"
            flags = re.IGNORECASE if params.get("ignore_case") else 0
            try:
                rx = re.compile(pattern, flags)
            except re.error as e:
                return f"正则错误: {e}"
            g = params.get("glob")
            max_results = int(params.get("max_results", 50))
            found = []
            try:
                it = root.rglob(g) if g else root.rglob("*")
            except Exception as e:
                return f"目录错误: {e}"
            for fp in it:
                if not fp.is_file():
                    continue
                if fp.stat().st_size > 5_000_000:
                    continue
                try:
                    lines = fp.read_text(encoding="utf-8", errors="ignore").splitlines()
                except Exception:
                    continue
                for i, line in enumerate(lines, 1):
                    if rx.search(line):
                        found.append(f"{fp}:{i}: {line.strip()[:200]}")
                        if len(found) >= max_results:
                            break
                if len(found) >= max_results:
                    break
            if not found:
                return f"未找到匹配: {pattern}"
            return f"匹配 {len(found)} 处:\n" + "\n".join(found)

        async def web_fetch(**params):
            """抓取网页 URL 并提取可读正文文本（对标 workbuddy 的网页研究能力）。
            params:
              url: 目标网址
              timeout: 可选，默认 20s
            返回纯文本正文（已去除脚本/样式/标签噪声），失败返回错误信息。
            """
            import urllib.request
            import urllib.error
            import html as _html
            import re as _re
            url = params.get("url", "")
            if not url:
                return "错误: 必须提供 url"
            timeout = float(params.get("timeout", 20))
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
            try:
                req = urllib.request.Request(url, headers=headers)
                with urllib.request.urlopen(req, timeout=timeout) as resp:
                    raw = resp.read()
                    ctype = resp.headers.get("Content-Type", "")
                    enc = "utf-8"
                    m = _re.search(r"charset=([\w-]+)", ctype)
                    if m:
                        enc = m.group(1)
                    body = raw.decode(enc, errors="ignore")
            except (urllib.error.URLError, OSError) as e:
                return f"抓取失败: {e}"
            # 去 <script>/<style>/<head>，再剥标签
            body = _re.sub(r"(?is)<script.*?</script>", " ", body)
            body = _re.sub(r"(?is)<style.*?</style>", " ", body)
            body = _re.sub(r"(?is)<head.*?</head>", " ", body)
            body = _re.sub(r"(?s)<[^>]+>", " ", body)
            body = _html.unescape(body)
            body = _re.sub(r"[ \t]+", " ", body)
            body = _re.sub(r"\n\s*\n+", "\n", body)
            text = body.strip()
            if len(text) > 20000:
                text = text[:20000] + "\n...[已截断]"
            return f"[来源 {url}]\n{text}"

        async def glob(**params):
            """按 glob 模式查找文件（对标 OpenCode glob，结果按修改时间排序）。
            params:
              pattern: 匹配模式，如 '*.py'、'**/*.py'、'src/**/*.ts'
              path: 搜索根目录，默认工作目录
              max_results: 最多返回数量，默认 200
            """
            import fnmatch
            pattern = params.get("pattern", "")
            if not pattern:
                return "错误: 必须提供 pattern"
            root = self._safe_path(Path(params.get("path", ".")))
            if not root.exists():
                return f"错误: 目录不存在 {root}"
            results = []
            try:
                for p in root.rglob("*"):
                    try:
                        rel = str(p.relative_to(root))
                    except Exception:
                        rel = str(p)
                    name = p.name
                    if (fnmatch.fnmatch(name, pattern)
                            or fnmatch.fnmatch(rel, pattern)
                            or fnmatch.fnmatch(str(p), pattern)):
                        try:
                            mtime = p.stat().st_mtime
                        except Exception:
                            mtime = 0
                        results.append((mtime, str(p)))
            except Exception as e:
                return f"glob 错误: {e}"
            results.sort(key=lambda x: -x[0])
            out = [f for _, f in results[:int(params.get("max_results", 200))]]
            if not out:
                return f"未匹配: {pattern}"
            return f"匹配 {len(out)} 个文件:\n" + "\n".join(out)

        async def websearch(**params):
            """网络搜索（对标 OpenCode websearch）。
            默认使用 DuckDuckGo HTML 解析；可在 agent_core/search_providers.json
            配置其它 provider（如带 API key 的搜索引擎）。
            params:
              query: 搜索词
              max_results: 最多结果，默认 8
              provider: 可选，覆盖默认 provider
            """
            query = params.get("query") or params.get("q") or ""
            if not query:
                return "错误: 必须提供 query"
            max_results = int(params.get("max_results", 8))
            import urllib.request, urllib.parse, urllib.error
            import html as _html, re as _re, json as _json
            # 可选 provider 配置
            provider = (params.get("provider") or "").strip()
            cfg_path = None
            try:
                from agent_core import xrz_paths
                cfg_path = xrz_paths.DATA_ROOT / "search_providers.json"
            except Exception:
                pass
            if cfg_path and cfg_path.exists():
                try:
                    provs = _json.loads(cfg_path.read_text(encoding="utf-8"))
                    prov = provs.get(provider or provs.get("default", ""), {})
                    if prov.get("url") and prov.get("parse") == "json":
                        return _websearch_json(query, prov, max_results)
                except Exception:
                    pass
            # 默认：DuckDuckGo HTML
            try:
                url = "https://html.duckduckgo.com/html/?q=" + urllib.parse.quote(query)
                req = urllib.request.Request(url, headers={
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"})
                with urllib.request.urlopen(req, timeout=float(params.get("timeout", 15))) as r:
                    page = r.read().decode("utf-8", "ignore")
                items = []
                for m in _re.finditer(
                        r'<a[^>]+class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>', page, _re.S):
                    link = m.group(1)
                    title = _re.sub(r"<[^>]+>", "", m.group(2)).strip()
                    link = _decode_ddg(link)
                    items.append({"title": title, "link": link})
                snips = _re.findall(r'class="result__snippet"[^>]*>(.*?)</a>', page, _re.S)
                for i, it in enumerate(items[:max_results]):
                    if i < len(snips):
                        it["snippet"] = _re.sub(r"<[^>]+>", "", snips[i]).strip()
                if not items:
                    return f"未找到搜索结果: {query}"
                lines = [f"搜索「{query}」共 {len(items)} 条:"]
                for i, it in enumerate(items[:max_results], 1):
                    lines.append(f"{i}. {it.get('title','')}\n   {it.get('link','')}\n   {it.get('snippet','')}")
                return "\n".join(lines)
            except (urllib.error.URLError, OSError) as e:
                return f"搜索失败: {e}（可检查网络或在 search_providers.json 配置带 API key 的 provider）"

        def _decode_ddg(link):
            """解码 DuckDuckGo 重定向链接中的真实 URL。"""
            import urllib.parse as _up
            m = _re.search(r"uddg=([^&]+)", link)
            if m:
                return _up.unquote(m.group(1))
            return link

        def _websearch_json(query, prov, max_results):
            import urllib.request, urllib.parse, urllib.error, json as _json, re as _re
            url = prov["url"] + urllib.parse.quote(query)
            headers = {"User-Agent": "Mozilla/5.0"}
            for k, v in (prov.get("headers") or {}).items():
                headers[k] = v
            try:
                req = urllib.request.Request(url, headers=headers)
                with urllib.request.urlopen(req, timeout=15) as r:
                    data = _json.loads(r.read().decode("utf-8", "ignore"))
                items = []
                for it in data[:max_results]:
                    items.append({
                        "title": str(it.get("title", "")),
                        "link": str(it.get("url") or it.get("link", "")),
                        "snippet": _re.sub(r"<[^>]+>", "", str(it.get("snippet") or it.get("description") or "")).strip(),
                    })
                if not items:
                    return f"未找到搜索结果: {query}"
                lines = [f"搜索「{query}」共 {len(items)} 条:"]
                for i, it in enumerate(items, 1):
                    lines.append(f"{i}. {it.get('title','')}\n   {it.get('link','')}\n   {it.get('snippet','')}")
                return "\n".join(lines)
            except Exception as e:
                return f"搜索失败: {e}"

        async def apply_patch(**params):
            """应用统一 diff 补丁（对标 opencode/codex 的 patch 应用）。
            params:
              patch: 统一 diff 文本（支持 --- a/.. +++ b/.. 前缀与 /dev/null 新建/删除文件）
              base_dir: 可选，补丁里相对路径的基准目录（默认工作目录）
            返回每个文件的应用结果摘要；冲突会逐文件报错而不中断其他文件。
            """
            import re as _re
            patch = params.get("patch", "")
            if not patch.strip():
                return "错误: 必须提供 patch"
            base = self._safe_path(Path(params.get("base_dir", ".")))

            def _norm(p):
                p = (p or "").strip()
                p = _re.sub(r"\t.*$", "", p)        # 去掉\t时间戳
                if p.startswith("a/") or p.startswith("b/"):
                    p = p[2:]
                if p in ("/dev/null", ""):
                    return None
                return p

            # 解析补丁为 文件 -> hunks
            files = []
            cur = None
            for raw in patch.split("\n"):
                if raw.startswith("--- "):
                    if cur:
                        files.append(cur)
                    cur = {"old": raw[4:], "new": None, "hunks": []}
                    continue
                if raw.startswith("+++ "):
                    if cur is not None:
                        cur["new"] = raw[4:]
                    continue
                if raw.startswith("@@"):
                    if cur is None:
                        continue
                    m = _re.search(r"@@ -(\d+)(?:,(\d+))? \+(\d+)(?:,(\d+))? @@", raw)
                    if not m:
                        continue
                    cur["hunks"].append({
                        "old_start": int(m.group(1)),
                        "new_start": int(m.group(3)),
                        "lines": [],
                    })
                    continue
                if raw.startswith("\\"):   # "\ No newline at end of file"
                    continue
                if cur is None or not cur["hunks"]:
                    continue
                if raw.startswith("+"):
                    cur["hunks"][-1]["lines"].append(("+", raw[1:]))
                elif raw.startswith("-"):
                    cur["hunks"][-1]["lines"].append(("-", raw[1:]))
                else:
                    cur["hunks"][-1]["lines"].append((" ", raw[1:] if raw.startswith(" ") else raw))
            if cur:
                files.append(cur)

            def _apply(f):
                old = _norm(f["old"]); new = _norm(f["new"])
                if old is None and new is None:
                    return "跳过: 路径无法解析"
                target_rel = new or old
                target = self._safe_path(base / target_rel)
                if new is None:   # 删除文件
                    if target.exists():
                        target.unlink()
                        return f"已删除: {target}"
                    return f"删除跳过(不存在): {target}"
                is_new = (old is None)
                if is_new:
                    # 新文件：原文件不存在，所有 + 行即为内容（忽略上下文/位置）
                    out = []
                    for h in f["hunks"]:
                        for tag, content in h["lines"]:
                            if tag == "+":
                                out.append(content)
                    target.parent.mkdir(parents=True, exist_ok=True)
                    target.write_text("\n".join(out), encoding="utf-8")
                    return f"已应用(新建): {target} (行数 0→{len(out)})"
                original = target.read_text(encoding="utf-8").split("\n") if target.exists() else []
                out = []
                oi = 0
                ok = True
                for h in f["hunks"]:
                    start = h["old_start"] - 1
                    if start > oi:
                        out.extend(original[oi:start]); oi = start
                    elif start < oi:
                        ok = False; break
                    for tag, content in h["lines"]:
                        if tag == " ":
                            out.append(content); oi += 1
                        elif tag == "-":
                            oi += 1
                        else:
                            out.append(content)
                if not ok:
                    return f"补丁冲突(行号错位): {target}"
                out.extend(original[oi:])
                target.parent.mkdir(parents=True, exist_ok=True)
                target.write_text("\n".join(out), encoding="utf-8")
                return f"已应用: {target} (行数 {len(original)}→{len(out)})"

            results = [_apply(f) for f in files]
            return "\n".join(results) if results else "错误: 补丁为空或格式无法识别"

        self._tools.register("file_edit", "编辑文件(replace/append/insert/delete, 对标OpenCode edit)", file_edit)
        self._tools.register("edit", "编辑文件(OpenCode 规范名, 同 file_edit)", file_edit)
        self._tools.register("grep", "按正则搜索文件内容(代码搜索, 对标OpenCode grep)", grep)
        self._tools.register("glob", "按模式查找文件(对标OpenCode glob, 如 **/*.py)", glob)
        self._tools.register("web_fetch", "抓取网页URL返回可读正文(对标OpenCode webfetch)", web_fetch)
        self._tools.register("webfetch", "抓取网页URL返回可读正文(OpenCode 规范名, 同 web_fetch)", web_fetch)
        self._tools.register("websearch", "网络搜索(DuckDuckGo, 对标OpenCode websearch)", websearch)
        self._tools.register("apply_patch", "应用统一diff补丁(对标OpenCode/Codex patch)", apply_patch)
        self._tools.register("patch", "应用统一diff补丁(OpenCode 规范名, 同 apply_patch)", apply_patch)

        # ─── 对标 OpenCode 的任务/权限/会话/MCP 框架 ───
        from agent_core import xrz_paths as _xz

        async def todowrite(**params):
            """任务列表管理（对标 OpenCode todowrite）。
            params（三选一）:
              todos: 任务数组整体替换（元素可为字符串，或 {content,status,activeForm}）
              content: 新增一条任务
              id: 更新指定任务（配合 content/activeForm/status）
            持久化到 DATA_ROOT/todos.json。"""
            import json as _json
            p = _xz.DATA_ROOT / "todos.json"
            todos = []
            if p.exists():
                try:
                    todos = _json.loads(p.read_text(encoding="utf-8"))
                except Exception:
                    todos = []
            if "todos" in params and isinstance(params["todos"], list):
                new = []
                for i, t in enumerate(params["todos"], 1):
                    if isinstance(t, str):
                        new.append({"id": i, "content": t, "status": "pending", "activeForm": t})
                    elif isinstance(t, dict):
                        new.append({"id": t.get("id", i), "content": t.get("content", ""),
                                    "status": t.get("status", "pending"),
                                    "activeForm": t.get("activeForm", t.get("content", ""))})
                todos = new
            elif params.get("content"):
                nid = max([t.get("id", 0) for t in todos], default=0) + 1
                todos.append({"id": nid, "content": params["content"], "status": "pending",
                              "activeForm": params.get("activeForm", params["content"])})
            elif params.get("id") is not None:
                tid = params["id"]
                for t in todos:
                    if t.get("id") == tid:
                        if "content" in params:
                            t["content"] = params["content"]
                        if "activeForm" in params:
                            t["activeForm"] = params["activeForm"]
                        if "status" in params:
                            t["status"] = params["status"]
                        break
            else:
                return "错误: 需提供 todos(整体替换) / content(新增) / id(更新)"
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(_json.dumps(todos, ensure_ascii=False, indent=2), encoding="utf-8")
            return _format_todos(todos)

        def _format_todos(todos):
            if not todos:
                return "任务列表为空（用 todowrite 添加任务）"
            done = sum(1 for t in todos if t.get("status") == "completed")
            lines = [f"任务列表 ({done}/{len(todos)} 完成):"]
            for t in todos:
                mark = {"completed": "✅", "in_progress": "🔄", "pending": "⬜"}.get(t.get("status", "pending"), "⬜")
                lines.append(f"  {mark} [{t.get('id')}] {t.get('content')}")
            return "\n".join(lines)

        async def todoread(**params):
            """读取当前任务列表（对标 OpenCode todoread）。"""
            import json as _json
            p = _xz.DATA_ROOT / "todos.json"
            if not p.exists():
                return _format_todos([])
            try:
                todos = _json.loads(p.read_text(encoding="utf-8"))
            except Exception:
                todos = []
            return _format_todos(todos)

        def _perm_path():
            return _xz.DATA_ROOT / "permissions.json"

        def _load_perms():
            import json as _json
            p = _perm_path()
            if p.exists():
                try:
                    return _json.loads(p.read_text(encoding="utf-8"))
                except Exception:
                    pass
            return {"default": "allow", "rules": []}

        async def permission_check(**params):
            """查询某工具/命令的权限（对标 OpenCode permission，支持通配符）。
            返回 {key, action}，action ∈ allow/deny/ask；未命中规则则用 default。"""
            import json as _json, fnmatch
            key = params.get("key") or params.get("tool") or ""
            if not key:
                return "错误: 必须提供 key"
            cfg = _load_perms()
            action = cfg.get("default", "allow")
            for rule in cfg.get("rules", []):
                rk = rule.get("key", "")
                if rk == key or fnmatch.fnmatch(key, rk):
                    action = rule.get("action", action)
            return _json.dumps({"key": key, "action": action}, ensure_ascii=False)

        async def permission_list(**params):
            """列出全部权限规则（对标 OpenCode permission）。"""
            import json as _json
            return _json.dumps(_load_perms(), ensure_ascii=False, indent=2)

        async def permission_set(**params):
            """设置一条权限规则（allow/deny/ask，支持通配符如 mcp_*、git *）。"""
            import json as _json
            key = params.get("key") or params.get("tool")
            action = params.get("action", "allow")
            if not key:
                return "错误: 必须提供 key"
            if action not in ("allow", "deny", "ask"):
                return "错误: action 必须是 allow/deny/ask"
            cfg = _load_perms()
            rules = cfg.setdefault("rules", [])
            found = False
            for r in rules:
                if r.get("key") == key:
                    r["action"] = action
                    found = True
                    break
            if not found:
                rules.append({"key": key, "action": action})
            p = _perm_path()
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(_json.dumps(cfg, ensure_ascii=False, indent=2), encoding="utf-8")
            return f"已设置权限: {key} -> {action}"

        def _save_sessions(p, sessions):
            import json as _json
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(_json.dumps(sessions, ensure_ascii=False, indent=2), encoding="utf-8")

        async def session_manager(**params):
            """多会话并行管理框架（对标 OpenCode 的 session/task 多会话能力）。
            action: list | current | create | switch
              create: name/platform/cwd 可选
              switch: id 指定会话
            会话元数据持久化到 DATA_ROOT/sessions.json，并广播 SESSION_UPDATED 事件。"""
            import json as _json
            action = (params.get("action") or "list").lower()
            p = _xz.DATA_ROOT / "sessions.json"
            sessions = []
            if p.exists():
                try:
                    sessions = _json.loads(p.read_text(encoding="utf-8"))
                except Exception:
                    sessions = []
            if action == "list":
                return _json.dumps({"sessions": sessions, "current": getattr(self, "_active_session_id", None)},
                                   ensure_ascii=False, indent=2)
            if action == "current":
                return _json.dumps({"current": getattr(self, "_active_session_id", None), "sessions": sessions},
                                   ensure_ascii=False)
            if action == "create":
                nid = max([s.get("id", 0) for s in sessions], default=0) + 1
                rec = {"id": nid, "name": params.get("name") or f"会话 {nid}",
                       "platform": params.get("platform", "deepseek"), "cwd": params.get("cwd", ""),
                       "created": time.time()}
                sessions.append(rec)
                self._active_session_id = nid
                _save_sessions(p, sessions)
                self._emit(EventType.SESSION_UPDATED, {"sessions": sessions, "current": nid})
                return _json.dumps(rec, ensure_ascii=False)
            if action == "switch":
                nid = params.get("id")
                for s in sessions:
                    if s.get("id") == nid:
                        self._active_session_id = nid
                        _save_sessions(p, sessions)
                        self._emit(EventType.SESSION_UPDATED, {"sessions": sessions, "current": nid})
                        return f"已切换到会话 {nid}: {s.get('name')}"
                return f"错误: 未找到会话 {nid}"
            return "错误: 未知 action(支持 list/current/create/switch)"

        async def task(**params):
            """委派子代理执行任务（对标 OpenCode task）。
            params: goal(必填) / type(research|visit|coding)"""
            goal = params.get("goal") or params.get("prompt") or ""
            if not goal:
                return "错误: 必须提供 goal"
            try:
                from agent_core.subagent_manager import get_subagent_manager
                mgr = get_subagent_manager(self._work_dir)
                tid = await mgr.spawn_subagent(goal, task_type=params.get("type", "research"))
                return f"已委派子代理任务: id={tid}\n目标: {goal}\n（用 check_task / wait_task 查询结果）"
            except Exception as e:
                return f"委派子代理失败: {e}"

        async def mcp_client(**params):
            """MCP 协议客户端框架（stdio 传输，对标 OpenCode 的 MCP 工具）。
            action:
              list_servers                      列出 mcp.json 中配置的 server
              list_tools  server=<名>           启动 server 并列出其 tools
              call_tool   server=<名> tool=<名> args=<obj>  调用某 tool
            配置：DATA_ROOT/mcp.json -> {"servers": {"名": {"command","args","env"}}}"""
            import json as _json
            action = (params.get("action") or "list_servers").lower()
            server = params.get("server", "")
            cfg_path = _xz.DATA_ROOT / "mcp.json"
            if not cfg_path.exists():
                return ("未配置 MCP：请在 " + str(cfg_path) + " 配置 servers，如：\n"
                        '{"servers": {"my": {"command": "npx", "args": ["-y", '
                        '"@modelcontextprotocol/server-everything"]}}}')
            try:
                cfgs = _json.loads(cfg_path.read_text(encoding="utf-8"))
            except Exception as e:
                return f"mcp.json 解析失败: {e}"
            servers = cfgs.get("servers", {})
            if action == "list_servers":
                return _json.dumps({"servers": list(servers.keys())}, ensure_ascii=False)
            if server not in servers:
                return f"未知 MCP server: {server}（可用: {list(servers.keys())}）"
            spec = servers[server]
            try:
                proc = await asyncio.create_subprocess_exec(
                    spec["command"], *spec.get("args", []),
                    stdin=asyncio.subprocess.PIPE, stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    env={**os.environ, **(spec.get("env") or {})})
            except Exception as e:
                return f"启动 MCP server 失败: {e}"
            counter = [0]

            async def _rpc(method, pms=None, notif=False):
                counter[0] += 1
                i = counter[0]
                msg = _json.dumps({"jsonrpc": "2.0", "id": i, "method": method, "params": pms or {}})
                proc.stdin.write((msg + "\n").encode())
                await proc.stdin.drain()
                if notif:
                    return None
                while True:
                    line = await proc.stdout.readline()
                    if not line:
                        return {"error": {"message": "MCP server 关闭连接"}}
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        r = _json.loads(line)
                    except Exception:
                        continue
                    if r.get("id") == i:
                        return r

            try:
                init = await _rpc("initialize", {"protocolVersion": "2024-11-05",
                                                 "capabilities": {}, "clientInfo": {"name": "xrz", "version": "1.0"}})
                await _rpc("notifications/initialized", {}, notif=True)
                if action == "list_tools":
                    res = await _rpc("tools/list", {})
                    tools = (res or {}).get("result", {}).get("tools", [])
                    return _json.dumps({"server": server, "tools": tools}, ensure_ascii=False)
                if action == "call_tool":
                    tname = params.get("tool", "")
                    targs = params.get("args") or {}
                    res = await _rpc("tools/call", {"name": tname, "arguments": targs})
                    return _json.dumps({"server": server, "tool": tname,
                                        "result": (res or {}).get("result")}, ensure_ascii=False)
                return _json.dumps({"server": server,
                                    "initialize": (init or {}).get("result")}, ensure_ascii=False)
            finally:
                try:
                    proc.terminate()
                except Exception:
                    pass

        self._tools.register("todowrite", "任务列表创建/更新/进度(对标OpenCode todowrite)", todowrite)
        self._tools.register("todoread", "读取任务列表(对标OpenCode todoread)", todoread)
        # 注：ask / question 工具已【刻意禁用】——agent 不得把任务踢回给用户，遇到不确定
        # 必须自己选合理默认方案用真实工具推进到底（见 _execute_command 中的拦截逻辑）。
        self._tools.register("permission_check", "权限查询(allow/deny/ask+通配符)", permission_check)
        self._tools.register("permission_list", "列出权限规则", permission_list)
        self._tools.register("permission_set", "设置权限规则(allow/deny/ask+通配符)", permission_set)
        self._tools.register("session_manager", "多会话并行管理框架", session_manager)
        self._tools.register("task", "委派子代理执行任务(对标OpenCode task)", task)
        self._tools.register("mcp_client", "MCP协议客户端框架(stdio)", mcp_client)

        # ─── 深度思考开关 + 文件附件工具 ───
        async def set_deep_think(**params):
            """开启/关闭当前 AI 平台的「深度思考」模式（每个平台 UI 不同，已逐平台适配）"""
            enable = params.get("enable", True)
            if self._session is None:
                return "错误：session 未初始化"
            try:
                await self._session.set_deep_think(enable=bool(enable))
                return f"深度思考已{'开启' if enable else '关闭'}"
            except Exception as e:
                return f"深度思考切换失败: {e}"

        async def attach_file(**params):
            """把本地文件（图片/PDF 等）附加到下次发送给 AI 平台的消息里。

            params.paths: 文件绝对/相对路径，支持单路径或路径列表。
            文件会被暂存，在下一轮对话时自动上传到输入框（出现预览/缩略图）。
            """
            paths = params.get("paths", [])
            if isinstance(paths, str):
                paths = [paths]
            if not paths:
                return "错误：paths 为空"
            from pathlib import Path as _P
            ok, missing = [], []
            for p in paths:
                if _P(p).exists():
                    ok.append(str(_P(p)))
                else:
                    missing.append(p)
            if not ok:
                return f"错误：所有文件都不存在：{missing}"
            self._pending_attachments.extend(ok)
            msg = f"已就绪附件 {len(ok)} 个"
            if missing:
                msg += f"（{len(missing)} 个不存在已跳过：{missing}）"
            msg += "，将在下一轮对话时自动附带发送给 AI"
            return msg

        self._tools.register("set_deep_think", "开/关深度思考模式(enable=true/false)", set_deep_think)
        self._tools.register("attach_file", "附加文件(图片/PDF等)到下次对话(支持paths列表)", attach_file)

        # ─── Shell 执行工具 ───
        async def shell_exec(**params):
            import asyncio
            cmd = params.get("command", "")
            # 注意：不再硬拦截「写 C 盘」的命令行——用户若明确要求在某个路径生成，
            # 就按用户路径执行，真实系统错误（如权限/路径不存在）会原样返回，
            # 由 AI 自行判断并纠正，脚本不替用户做路径决定。
            timeout = params.get("timeout", 60)
            try:
                proc = await asyncio.create_subprocess_shell(
                    cmd,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                )
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
                return f"[退出码 {proc.returncode}]\n{stdout.decode('utf-8', errors='ignore')}\n{stderr.decode('utf-8', errors='ignore')}"
            except asyncio.TimeoutError:
                proc.kill()
                return f"命令超时（>{timeout}s）"
            except Exception as e:
                return f"执行错误: {e}"

        self._tools.register("shell_exec", "执行Shell命令", shell_exec)
        self._tools.register("bash", "执行Shell命令(OpenCode 规范名, 同 shell_exec)", shell_exec)

        # ─── 浏览器操作工具（母代理直接执行）───

        async def browser_click(**params):
            return await self._bm.browser_click(params.get("selector", ""))

        async def browser_fill(**params):
            return await self._bm.browser_fill(params.get("selector", ""), params.get("text", ""))

        async def browser_screenshot(**params):
            path = params.get("path", "screenshot.png")
            full_path = Path(self._work_dir) / path
            return await self._bm.browser_screenshot(str(full_path))

        async def browser_search(**params):
            """网页搜索工具 - 不使用浏览器，直接用 HTTP 请求抓取页面
            
            通过 Bing 搜索获取结果，抓取页面内容后返回给 AI。
            不会占用母代理的浏览器状态。
            """
            import subprocess
            import json
            from pathlib import Path
            
            query = params.get("query", "")
            max_pages = params.get("max_pages", 3)
            output_file = params.get("output_file", "")
            
            if not query:
                return "[错误] 缺少 query 参数"
            
            # 确定输出文件路径
            if not output_file:
                output_file = str(Path(self._work_dir) / "search_results.json")
            
            # 调用网页搜索脚本
            script_path = Path(__file__).parent.parent / "web_searcher.py"
            
            try:
                result = subprocess.run(
                    [sys.executable, str(script_path), query, str(max_pages), output_file],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                
                if result.returncode == 0:
                    # 读取结果文件
                    if Path(output_file).exists():
                        with open(output_file, "r", encoding="utf-8") as f:
                            data = json.load(f)
                        return f"[搜索完成] 抓取 {data.get('scraped_count', 0)} 个页面\n\n{data.get('findings', '')}"
                    else:
                        return result.stdout
                else:
                    return f"[搜索错误] {result.stderr}"
                    
            except subprocess.TimeoutExpired:
                return "[错误] 搜索超时（60秒）"
            except Exception as e:
                return f"[错误] {str(e)}"

        # ─── 深度思考工具（母代理执行，AI 调用）───
        async def _deep_think_tool(**params):
            """深度思考工具 - 由 AI 调用，实际控制浏览器按钮
            
            当 AI 认为当前问题需要深度思考时，调用此工具开启 DeepSeek 的深度思考模式。
            系统会在 AI 输出完成后自动关闭此模式。
            """
            enable_str = params.get("enable", "true")
            enable = enable_str.lower() in ("true", "1", "yes", "on")
            await self._session.set_deep_think(enable)
            return f"深度思考模式已{'开启' if enable else '关闭'}"

        self._tools.register("browser_click", "点击元素", browser_click)
        self._tools.register("browser_fill", "填写输入框", browser_fill)
        self._tools.register("browser_screenshot", "截图", browser_screenshot)
        self._tools.register("browser_search", "搜索", browser_search)
        self._tools.register("deep_think", "深度思考：AI认为需要深度思考时调用此指令开启DeepSeek深度思考模式，完成后自动关闭", _deep_think_tool)

        # ─── 研究代理工具（独立子代理进程）───
        async def _browser_research(**params):
            """启动研究子代理（同一浏览器、不同窗口，共享登录态）"""
            from .subagent_manager import get_subagent_manager
            sam = get_subagent_manager(self._work_dir)
            # 把母代理的浏览器管理器注入，子代理将在同一浏览器里开新窗口
            sam.set_browser_manager(self._bm)

            query = params.get("query", "")
            max_pages = params.get("max_pages", 5)

            # 构建完整查询
            full_query = f"深度研究: {query} (最多访问 {max_pages} 个页面)"

            # 启动子代理（非阻塞，进程内、同浏览器）
            task_id = await sam.spawn_subagent(full_query, task_type="research")

            return f"[子代理已启动] task_id={task_id}\n查询: {query[:50]}...\n使用 check_task('{task_id}') 查看进度"

        async def _browser_visit(**params):
            """启动访问子代理（同一浏览器、不同窗口，共享登录态）"""
            from .subagent_manager import get_subagent_manager
            sam = get_subagent_manager(self._work_dir)
            # 把母代理的浏览器管理器注入，子代理将在同一浏览器里开新窗口
            sam.set_browser_manager(self._bm)

            url = params.get("url", "")

            # 构建完整查询
            full_query = f"访问并分析网页: {url}"

            # 启动子代理（非阻塞，进程内、同浏览器）
            task_id = await sam.spawn_subagent(full_query, task_type="visit")

            return f"[子代理已启动] task_id={task_id}\nURL: {url}\n使用 check_task('{task_id}') 查看进度"
        
        async def _check_task(**params):
            """检查子代理任务状态（含详细结果）"""
            from .subagent_manager import get_subagent_manager
            sam = get_subagent_manager(self._work_dir)

            task_id = params.get("task_id", "")
            task = sam.check_task(task_id)

            if not task:
                return f"任务 {task_id} 不存在"

            status_info = (
                f"任务: {task_id}\n"
                f"状态: {task.status.value}\n"
                f"类型: {task.task_type}\n"
                f"查询: {task.query[:80]}"
            )

            if task.result:
                r = task.result
                status_info += f"\n\n=== 结果 ===\n成功: {r.success}"
                if r.error:
                    status_info += f"\n错误: {r.error}"
                # findings：子代理的核心发现/产出内容（最重要）
                if r.findings:
                    findings_preview = r.findings[:1000] + "..." if len(r.findings) > 1000 else r.findings
                    status_info += f"\n\n--- 核心发现 ---\n{findings_preview}"
                # output：子代理最终回复文本
                if r.output:
                    output_preview = r.output[:500] + "..." if len(r.output) > 500 else r.output
                    status_info += f"\n\n--- 回复文本 ---\n{output_preview}"
                # 文件列表（带路径和大小）
                if r.files:
                    status_info += f"\n\n--- 产物文件 ({len(r.files)} 个) ---"
                    for fi in r.files[:10]:
                        name = fi.get("name", fi.get("path", "?"))
                        size = fi.get("size", 0)
                        preview = fi.get("preview", "")[:120] if fi.get("preview") else ""
                        status_info += f"\n  [{name}] ({size} bytes)"
                        if preview:
                            status_info += f" 预览: {preview}"

            return status_info

        async def _wait_task(**params):
            """等待子代理任务完成（返回完整结果）"""
            from .subagent_manager import get_subagent_manager
            sam = get_subagent_manager(self._work_dir)

            task_id = params.get("task_id", "")
            timeout = params.get("timeout", 300)

            task = await sam.wait_task(task_id, timeout=timeout)

            if not task:
                return f"等待任务 {task_id} 超时或任务不存在"

            if task.result and task.result.success:
                r = task.result
                parts = [f"任务 {task_id} 完成！"]
                if r.findings:
                    findings_text = r.findings[:2000] + "..." if len(r.findings) > 2000 else r.findings
                    parts.append(f"\n=== 核心发现/产出 ===\n{findings_text}")
                if r.output:
                    output_text = r.output[:800] + "..." if len(r.output) > 800 else r.output
                    parts.append(f"\n=== 子代理回复 ===\n{output_text}")
                if r.files:
                    parts.append(f"\n=== 产物文件 ({len(r.files)} 个) ===")
                    for fi in r.files[:15]:
                        name = fi.get("name", fi.get("path", "?"))
                        parts.append(f"  - {name} ({fi.get('size', 0)} bytes)")
                return "\n".join(parts)
            else:
                error = task.result.error if task.result else "未知错误"
                return f"任务 {task_id} 失败: {error}"

        self._tools.register("browser_research", "启动研究子代理（独立进程）", _browser_research)
        self._tools.register("browser_visit", "启动访问子代理（独立进程）", _browser_visit)
        self._tools.register("check_task", "检查子代理任务状态", _check_task)
        self._tools.register("wait_task", "等待子代理任务完成", _wait_task)

        # ─── Word 文档生成工具（母代理直接执行）───
        async def docx_tool_fn(**params):
            import os
            from pathlib import Path

            raw_content = ""
            for key in ["content", "text", "body", "data", "value", "raw"]:
                if key in params and params[key] is not None:
                    raw_content = params[key]
                    break

            # 清洗 DeepSeek/通义等网页 UI 残留：代码围栏、『json 复制 下载』『复制 下载』等
            # 工具栏文字会混进模型输出的内容里，必须剥掉，否则生成的文档满是垃圾行。
            def _clean_docx_text(t: str) -> str:
                import re
                out = []
                for ln in str(t).splitlines():
                    s = ln.rstrip()
                    if s.strip().startswith("```"):
                        continue
                    s = re.sub(r"^json\s*复制\s*下载\s*$", "", s)
                    s = re.sub(r"复制\s*下载\s*$", "", s)
                    s = re.sub(r"^\s*下载\s*$", "", s)
                    if s.strip():
                        out.append(s)
                return "\n".join(out)

            if isinstance(raw_content, str):
                raw_content = _clean_docx_text(raw_content)

            requested_filename = params.get("filename")
            filename = os.path.expandvars(str(requested_filename or "report.docx"))
            path = os.path.expandvars(str(params.get("path") or ""))
            filename_path = Path(filename)
            if filename_path.suffix.lower() != ".docx":
                filename_path = filename_path.with_suffix(".docx")

            # AI 有时把完整文件名误放在 path 中；此时不能再拼 report.docx，
            # 否则会产生 xxx.docx/report.docx 这种错误目录。
            if path and Path(path).suffix.lower() == ".docx" and not requested_filename:
                out_path = self._safe_path(Path(path))
            elif path:
                out_path = self._safe_path(Path(path) / filename_path)
            else:
                out_path = self._safe_path(filename_path)
            out_path.parent.mkdir(parents=True, exist_ok=True)

            try:
                from docx import Document

                doc = Document()
                structured = isinstance(raw_content, dict) and isinstance(raw_content.get("sections"), list)
                if structured:
                    title = str(raw_content.get("title") or raw_content.get("filename") or out_path.stem)
                else:
                    # 非结构化：优先用内容里第一个 # 标题当文档标题，否则回退文件名
                    _first_h = ""
                    if isinstance(raw_content, str):
                        for _ln in raw_content.splitlines():
                            if _ln.strip().startswith("# "):
                                _first_h = _ln.strip()[2:].strip()
                                break
                    title = _first_h or out_path.stem
                doc.add_heading(title, 0)

                def add_content(value, *, list_style=None):
                    if value is None:
                        return
                    if isinstance(value, dict):
                        for key, item in value.items():
                            doc.add_heading(str(key), level=2)
                            add_content(item)
                        return
                    if isinstance(value, (list, tuple)):
                        for item in value:
                            add_content(item, list_style="List Bullet")
                        return
                    for raw_line in str(value).splitlines() or [""]:
                        line = raw_line.strip()
                        if not line:
                            continue
                        if line.startswith("### "):
                            doc.add_heading(line[4:], level=3)
                        elif line.startswith("## "):
                            doc.add_heading(line[3:], level=2)
                        elif line.startswith("# "):
                            doc.add_heading(line[2:], level=1)
                        elif line.startswith(("- ", "* ")):
                            doc.add_paragraph(line[2:], style="List Bullet")
                        else:
                            doc.add_paragraph(line, style=list_style)

                if structured:
                    for index, section in enumerate(raw_content["sections"], start=1):
                        if isinstance(section, dict):
                            heading = section.get("heading") or section.get("title") or f"第{index}节"
                            section_content = section.get("content", section.get("body", section.get("text", "")))
                        else:
                            heading = f"第{index}节"
                            section_content = section
                        doc.add_heading(str(heading), level=1)
                        add_content(section_content)
                elif isinstance(raw_content, dict):
                    add_content(raw_content)
                else:
                    add_content(raw_content)

                doc.save(str(out_path))
                return f"✅ Word 文档已生成: {out_path}"
            except ImportError as ie:
                raise ImportError("❌ python-docx 未安装。请运行：pip install python-docx") from ie
            except Exception as e:
                raise Exception(f"❌ Word 文档生成失败 ({type(e).__name__}): {e}") from e

        self._tools.register("docx_create", "生成 Word 文档（python-docx）", docx_tool_fn)

        # ─── PPT 生成工具 ───
        async def pptx_tool_fn(**params):
            import os
            from pathlib import Path

            raw_content = ""
            for key in ["content", "text", "body", "data", "value", "raw"]:
                if key in params and params[key] is not None:
                    raw_content = params[key]
                    break

            requested_filename = params.get("filename")
            filename = os.path.expandvars(str(requested_filename or "slides.pptx"))
            path = os.path.expandvars(str(params.get("path") or ""))
            filename_path = Path(filename)
            if filename_path.suffix.lower() != ".pptx":
                filename_path = filename_path.with_suffix(".pptx")

            if path and Path(path).suffix.lower() == ".pptx" and not requested_filename:
                out_path = self._safe_path(Path(path))
            elif path:
                out_path = self._safe_path(Path(path) / filename_path)
            else:
                out_path = self._safe_path(filename_path)
            out_path.parent.mkdir(parents=True, exist_ok=True)

            try:
                from pptx import Presentation

                prs = Presentation()
                structured = isinstance(raw_content, dict) and isinstance(raw_content.get("slides"), list)
                deck_title = (
                    str(raw_content.get("title") or raw_content.get("filename") or out_path.stem)
                    if structured else out_path.stem
                )

                def normalize_body(value):
                    if value is None:
                        return []
                    if isinstance(value, dict):
                        return [f"{key}: {item}" for key, item in value.items()]
                    if isinstance(value, (list, tuple)):
                        return [str(item) for item in value]
                    return [line.strip() for line in str(value).splitlines() if line.strip()]

                if structured:
                    slides_data = raw_content["slides"]
                    if not slides_data:
                        slides_data = [{"heading": deck_title, "content": ""}]
                    for index, slide_data in enumerate(slides_data, start=1):
                        if isinstance(slide_data, dict):
                            slide_title = slide_data.get("heading") or slide_data.get("title") or f"第{index}页"
                            slide_content = slide_data.get("content", slide_data.get("body", slide_data.get("text", "")))
                        else:
                            slide_title = f"第{index}页"
                            slide_content = slide_data
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = str(slide_title)
                        frame = slide.shapes.placeholders[1].text_frame
                        frame.clear()
                        lines = normalize_body(slide_content)
                        for line_index, line in enumerate(lines):
                            paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
                            paragraph.text = line
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = deck_title
                    frame = slide.shapes.placeholders[1].text_frame
                    frame.clear()
                    lines = normalize_body(raw_content)
                    for line_index, line in enumerate(lines):
                        paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
                        paragraph.text = line

                prs.save(str(out_path))
                return f"✅ PPT 已生成: {out_path}"
            except ImportError as ie:
                raise ImportError("❌ python-pptx 未安装。请运行：pip install python-pptx") from ie
            except Exception as e:
                raise Exception(f"❌ PPT 生成失败 ({type(e).__name__}): {e}") from e

        self._tools.register("pptx_create", "生成 PPT（python-pptx）", pptx_tool_fn)

        # ─── 记忆工具（母代理直接执行）───
        async def remember_tool(**params):
            if self._memory is None:
                return "错误：记忆管理器未初始化"
            content = params.get("content", "")
            tags = params.get("tags", [])
            mid = self._memory.save(content, tags=tags if isinstance(tags, list) else [tags])
            return f"[记忆已保存] ID={mid}"

        async def recall_tool(**params):
            if self._memory is None:
                return "错误：记忆管理器未初始化"
            query = params.get("query", "")
            results = self._memory.search(query)
            if not results:
                return "[回忆] 未找到相关记忆"
            lines = [f"- [{r.id}] {r.content[:60]}... (相关度: {r.score:.2f})" for r in results]
            return "[回忆结果]\n" + "\n".join(lines)

        async def summarize_tool(**params):
            if self._memory is None:
                return "错误：记忆管理器未初始化"
            period = params.get("period", "today")
            results = self._memory.summarize(period=period)
            return f"[摘要] 共 {len(results)} 条记忆"

        async def list_summaries_tool(**params):
            if self._memory is None:
                return "错误：记忆管理器未初始化"
            limit = params.get("limit", 10)
            return self._memory.list(limit=limit)

        self._tools.register("remember", "保存记忆", remember_tool)
        self._tools.register("recall", "回忆记忆", recall_tool)
        self._tools.register("summarize", "生成摘要", summarize_tool)
        self._tools.register("list_summaries", "列出记忆摘要", list_summaries_tool)

        # ─── 子代理结果查询工具 ───
        async def get_subagent_result_tool(**params):
            if self._subagent_manager is None:
                return "错误：子代理管理器未初始化"
            task_id = params.get("task_id", "")
            task = self._subagent_manager.check_task(task_id)
            if not task:
                return f"任务 {task_id} 不存在"
            if task.status.value != "done":
                return f"任务 {task_id} 状态: {task.status.value}（尚未完成）"
            r = task.result
            parts = [f"[子代理结果] 任务 {task_id}"]
            if r.findings:
                parts.append(f"\n=== 核心发现 ===\n{r.findings[:1500]}")
            if r.output:
                parts.append(f"\n=== 回复 ===\n{r.output[:500]}")
            if r.files:
                parts.append(f"\n=== 文件 ({len(r.files)} 个) ===")
                for fi in r.files[:10]:
                    parts.append(f"  - {fi.get('name', fi.get('path', '?'))}")
            return "\n".join(parts)

        self._tools.register("get_subagent_result", "获取子代理结果", get_subagent_result_tool)

        # ─── Ollama 本地模型对话工具 ───
        async def ollama_chat_tool(**params):
            """通过 Ollama REST API 与本地模型对话"""
            from agent_core.multi_browser import get_multi_browser_manager
            mgr = get_multi_browser_manager()
            if not mgr or not mgr.ollama:
                return "[错误] Ollama 未初始化。请确保 multi_browser 已启动"
            
            prompt = params.get("prompt", "")
            system = params.get("system", "")
            if not prompt:
                return "[错误] 缺少 prompt 参数"
            
            result = await mgr.ollama_chat(prompt, system)
            if result.get("success"):
                return f"[Ollama/{result.get('model', '?')}]\n{result['text']}"
            else:
                return f"[Ollama 错误] {result.get('error', '未知')}"

        self._tools.register("ollama_chat", "通过 Ollama 本地 API 对话（自动检测可用模型）", ollama_chat_tool)

    # ============================================================
    # 系统提示词
    # ============================================================

    def _build_system_prompt(self) -> str:
        """构建系统提示词（包含所有可用工具）"""
        tool_desc_list = []
        for name, info in self._tools._tools.items():
            tool_desc_list.append(f"- {name}: {info['description']}")
        tools_block = "\n".join(tool_desc_list)

        # 多平台 LLM 支持：从平台注册表动态生成（新增网页 AI 只需改 platforms.json）
        try:
            _plats = _list_platforms()
        except Exception:
            _plats = []
        if _plats:
            platforms_block = "\n".join(
                f"- {p['name']}（{p.get('chat_url') or p.get('url')}）" for p in _plats
            )
            platforms_block += "\n- Ollama（本地子进程: ollama run qwen3.5:0.8b，不需要浏览器）"
            platforms_block += "\n\n更多网页 AI 可通过编辑 agent_core/platforms.json 或 "
            platforms_block += "<数据目录>/platforms.user.json 即时接入，无需改动代码。"
        else:
            platforms_block = "- DeepSeek（https://chat.deepseek.com，默认平台）"

        return f"""你是仙人掌 Agent（XianRenZhang Agent），一个自主 AI 助手，支持多平台 LLM。

=== 核心指令 ===
当用户提出任务时，你必须根据需要使用工具来完成任务。工具调用格式：

@@@@
{{
    "tool": "工具名称",
    "params": {{...}},
    "id": "唯一标识"
}}
@@@@

RAW 命令格式（直接执行shell）：
<<<RAW>>>
命令内容
<<<RAW>>>

=== 可用工具 ===
{tools_block}

工具使用约定（对标 OpenCode / Codex）：
- 文件操作：read(支持 offset/limit 行范围) / write / edit（mode=replace|append|insert|delete，replace 支持 old/new 或 pattern/repl 正则）/ list / glob(**/*.py) / grep(正则) / apply_patch(统一diff，支持新建/修改/删除)。
- 路径：相对路径基于工作目录；绝对路径按用户要求原样使用（不擅自改写到 C 盘）。
- 网络：webfetch(抓网页正文) / websearch(网络搜索，默认 DuckDuckGo)。
- 任务编排：todowrite/todoread(任务列表与进度) / task(委派子代理并行执行) / session_manager(多会话并行管理)。
- 权限：permission_check/permission_list/permission_set（allow/deny/ask + 通配符，如 mcp_*、git *）。
- 扩展：mcp_client 可对接任意 MCP server（stdio，配置见 DATA_ROOT/mcp.json）。

=== 多平台 LLM 支持 ===
你可以使用以下平台（需要登录或本地运行）：
{platforms_block}

不同平台的特性：
- DeepSeek: 支持深度思考(R1)、文件上传、搜索 —— 用 set_deep_think(enable=true) 开启
- 通义千问: 支持深度思考、文件上传 —— 用 set_deep_think(enable=true) 开启
- 豆包: 支持深度思考、文件上传 —— 用 set_deep_think(enable=true) 开启
- 元宝: 支持深度思考、文件上传 —— 用 set_deep_think(enable=true) 开启
- ChatGPT: 支持深度思考（Analysis 模式）、文件上传
- Gemini: 支持深度思考（Think 模式）
- Ollama: 完全本地运行，通过 `ollama run qwen3.5:0.8b` 子进程通信。AI 调用 ollama_chat 工具时启动 ollama run 并与 qwen3.5:0.8b 对话。

提示：需要深度思考时，先调用 `set_deep_think(enable=true)`，再发后续对话；
需要让 AI 看图/读文档时，先 `attach_file(paths=[...])` 把图片/PDF 路径加入，
再发对话，文件会自动作为附件上传到输入框。

=== 关键规则 ===
1. **禁止直接在 DeepSeek 界面搜索** - 所有搜索必须使用 `browser_search` 工具
2. 一次只能执行一个工具调用
3. 工具执行后，系统会返回结果，你需要根据结果决定下一步
4. 如需多步骤，请分多次发送工具调用
5. 任务完成后，发送 done() 表示结束
7. **思考过程可视化**：每次调用工具前，先用一两句自然语言表达你的思考或计划，
   并【写在本轮 @@@@ 协议块之外】（纯自然语言，不要放进 JSON）。这部分会被实时抓取，
   作为「🧠 思考过程」展示给用户，让用户看到你的推理链路。即便用户没开启深度思考(R1)，
   这条思考过程也必须保留。

8. **文件类任务必须用工具真实完成，绝不许让用户复制粘贴**：你运行在用户的本机电脑上，
   `docx_create` / `pptx_create` / `file_write` / `file_edit` 等工具由本地运行时
   【真实执行】，会在用户电脑上真的生成或修改文件。当用户要求生成 Word / PPT / 文本文件、
   或在某文件末尾追加内容时，你必须调用对应工具并传入完整参数，【绝不要】输出
   「我无法在您的电脑上创建文件，请复制以下内容到 Word」之类的话——那是错误的，
   你其实可以创建。示例：
   - 生成 Word 报告到指定路径：
     @@@@{{"tool":"docx_create","params":{{"content":"# 仙人掌 Agent 自测报告\\n一、功能概览\\n...","path":"C:\\\\Users\\\\X.LAPTOP-CA1GJQE3\\\\Desktop\\\\test\\\\report.docx"}},"id":"1"}}@@@@
   - 在文件末尾追加一行：
     @@@@{{"tool":"file_edit","params":{{"path":"C:\\\\Users\\\\X.LAPTOP-CA1GJQE3\\\\Desktop\\\\test\\\\note.txt","mode":"append","new":"TEST_OK_仙人掌自测"}},"id":"2"}}@@@@
   - 注意工具名必须是上面「可用工具」里列的真名（docx_create / file_edit 等）；
     若你一时想不起准确名，本地运行时会自动把相近的名字映射到正确工具，但请尽量用真名。

=== 任务范围与「绝不偷懒」铁律（极重要）===
- 你【只能】执行用户明确要求的任务，以及完成该任务直接必需的子步骤。
- 【严禁】自行发明新任务；【严禁】主动生成用户没有要求的报告 / 调研 / 文件 / 网页 / 对话。
- 历史对话【仅作上下文参考】，绝不自动续做、绝不重发历史中未完成或未确认的任务。
  即使恢复了历史对话，也【绝不】把历史里任何一条用户消息当作新指令去执行。
- 每次用户发来一条新指令，你只处理那一条指令本身，不要把旧历史混进来自行扩展。
- 【最重要 / 绝不偷懒】你【没有】 ask / question 这类「向用户提问」的工具，也【绝不允许】
  以任何形式把任务踢回给用户、或假装「做不了 / 我无法」。遇到信息不全、路径不确定、
  或拿不准时，自己选择【最合理的默认方案】直接推进，用真实工具把用户的任务做到底，
  然后调用 done() 结束。
- 如果你觉得某件事「用户好像没明确要求」，但它是完成当前这条指令的【直接必需子步骤】
  （例如为了生成报告必须先建目录、为了改文件必须先读取），那就直接做，不要停下来。
- 只有当你想做的东西【完全超出】用户这条指令的范畴（用户没要报告你偏要写报告）时，才不做它；
  但即便如此也【不向用户提问】，只专注于用户真正要求的事，做完后调用 done()。

=== 子代理规则 ===
- browser_research: 启动独立研究子代理（后台运行，不阻塞母代理）
- browser_visit: 启动独立访问子代理
- 子代理启动后返回 task_id，使用 check_task(task_id) 查询进度
- 子代理不能再创建子代理（MAX_DEPTH=1）

=== 记忆系统 ===
- remember(content, tags): 保存重要信息到长期记忆
- recall(query): 搜索相关记忆
- 每运行约 20 轮会自动提醒保存记忆

=== 文件输出路径 ===
- 产物默认落在工作目录（D 盘项目目录）。
- 【关键】如果用户明确指定了输出路径，你必须**严格按用户指定的路径生成文件**，
  任何盘、任何目录都可以（包括 C 盘），不要擅自改动盘符或目录，也不要重定向到别处。
- 软件自身安装在 D 盘、其内部数据（登录态、对话历史、缓冲、子代理输出、Playwright
  浏览器二进制）由 xrz_paths.py 固定存放在 D 盘项目目录——这跟你「工作时生成产物的
  路径」是两回事，互不影响。你工作时没有盘符限制，用户叫你生成在哪就在哪。

=== 任务结束（由你判断）===
- 任务是否完成、**何时结束，由你自己判断**：工作确实做完了，就调用 done() 结束。
- 不要假称「失败 / 做不了」——如果工具已返回成功，就说明产物已经生成，如实收尾即可。
- 严禁在任务没做完时提前放弃；也没必要在成功后反复「补救」一个已经成功的产物。
- 你自己的上一轮回复【不会被回灌】给你（避免回声干扰）。如需回顾你做过什么，
  调用 recall("...") 从记忆里查询执行进度。

=== 会话历史追溯 ===
- 每次对话会生成唯一的 DeepSeek URL，可通过该 URL 追溯历史
- 使用 save_conversation() 保存当前对话上下文
- 使用 load_conversation(url=...) 加载历史对话

=== 工作目录 ===
所有文件操作默认相对于: {self._work_dir}
（此目录位于 D 盘，是默认落点；用户明确指定的其它路径同样可用，不限制盘符）
"""

    # ============================================================
    # 公共 API
    # ============================================================

    async def start(self, session: Optional[DeepSeekSession] = None):
        """启动 Commander（复用外部 session）"""
        if session:
            self._session = session
        self._running = True
        # 设置系统提示词到 session，这样每轮对话都会包含
        if self._session:
            self._session.set_system_prompt(self._system_prompt)

    async def run(self, user_instruction: str, file_path: Optional[str] = None,
                  context_hints: str = "") -> str:
        """
        单轮执行（自动循环直到 AI 认为完成）
        返回最终回复内容
        """
        if self._session is None:
            raise RuntimeError("Session 未初始化")

        # 【关键修复】每条命令都强制把系统提示词（工具列表 + @@@@ 协议）注入 session，
        # 否则新建/restore 的会话 self._system_prompt 为空，模型收不到工具说明会坚称
        # 「我无法操作你的电脑」，导致文件类任务全部失败。
        self._session.set_system_prompt(self._system_prompt)

        # 构建第一轮输入（包含系统提示词只在这一轮，后续已存入session）
        current_input = user_instruction
        original_task = user_instruction  # 任务范围守卫：用于拦截 AI 自行发明的新任务
        if file_path:
            current_input += f"\n\n[参考文件: {file_path}]"
        if context_hints:
            current_input += f"\n\n[上下文提示: {context_hints}]"

        final_reply = ""
        last_ai_text = ""  # 最近一次 AI 的纯文本回复（循环到上限时作为最终答案返回）
        # 「动作-结果」摘要：每轮记录，回灌给模型的是这个而不是 agent 原话
        action_log: list = []
        # 【硬性约束】循环只能通过 AI 显式调用 done() 并经脚本确认后退出。
        # max_turns 设为极大值（≈ 无穷大），绝不因为轮数限制提前踢出 AI。
        # 唯一合法退出路径：AI 发送 {"tool":"done"} → 脚本确认 → break。
        # 安全阀：MAX_NO_PROTOCOL 防止 AI 永远不回协议导致纯纠正死循环（独立于轮数）。
        max_turns, turn = 99999999999, 0
        no_protocol_retries = 0  # 防止通用平台（通义/豆包等）不遵守协议时无限循环
        # 极大值：通用平台不一定每次都立刻回 @@@@ 协议（尤其首轮还在消化协议时），
        # 绝不能因为几次没回协议就提前 break 退出循环（用户明确要求调成极大值）。
        # 配合「每轮反复发送系统提示词」+「清洗网页 UI 文字」，模型最终会回协议。
        MAX_NO_PROTOCOL = 99999999999

        while turn < max_turns and self._running:
            turn += 1
            self._history_turns += 1

            # 记忆提醒
            if self._memory and self._history_turns % 10 == 0:
                self._emit(EventType.MEMORY_REMINDER, {"turn": self._history_turns})

            # 发送给 AI（携带待附件）
            try:
                response = await self._session.send(
                    current_input,
                    attachments=self._pending_attachments if self._pending_attachments else None,
                )
                self._pending_attachments = []  # 发送后清空
            except Exception as e:
                logger.error(f"AI 调用失败: {e}")
                return f"[错误] AI 调用失败: {e}"

            # 诊断日志：记录收到的回复长度和是否包含协议标记
            has_protocol = "@@@@" in (response or "")
            logger.info(
                f"[Commander] 第 {turn} 轮收到回复: {len(response or '')} 字符, "
                f"含 @@@@ 协议: {has_protocol}"
            )

            # 解析指令
            cmds = self._protocol.extract_all(response)

            # ── 思考过程可视化 ──
            # 把 AI 在 @@@@ 协议 / <<<RAW>>> 之外的自然语言表达（思考、计划、解释）推送给 GUI，
            # 作为「🧠 思考过程」实时展示。这样即便平台没有开启深度思考(R1)推理模式，
            # 用户也能看到 Agent 每一步在想什么、打算怎么做。
            _reasoning = self._strip_protocol(response)
            if _reasoning:
                if not cmds:
                    # AI 这一轮只说了话、没有调工具 → 直接当作思考过程展示
                    self._emit(EventType.THINKING, {"text": _reasoning})
                elif cmds[0].command.tool != "done":
                    # 调工具前的计划/解释 → 展示为思考过程
                    # （done 轮的最终答复会由 ai_final_reply 单独展示，避免重复）
                    self._emit(EventType.THINKING, {"text": _reasoning})

            # 若脚本自动修复了 AI 的 JSON 语法错误，明确记录/告知（透明、便于排查）
            if cmds and getattr(cmds[0], "fixed", False):
                self._emit(EventType.CORRECTION_SENT,
                           {"text": self._protocol.describe_fix(cmds[0])})
            if not cmds:
                ai_text = (response or "").strip()
                if ai_text:
                    last_ai_text = ai_text  # 记住最近一次真实回答，兜底用
                    no_protocol_retries += 1
                    if no_protocol_retries <= MAX_NO_PROTOCOL:
                        current_input = (
                            "[SYSTEM] NO @@@@ PROTOCOL. Your text was: " + ai_text[:200] + ". All responses MUST use @@@@ JSON format. Format: @@@@\\n{\"tool\":\"xxx\",\"params\":{},\"id\":\"1\"}\\n@@@@. Please retry."
                        )
                        continue
                    else:
                        # 通用平台（通义/豆包/元宝）不一定会输出 @@@@ 协议，
                        # 多次纠正无效后，把这段文本当作最终回复返回，避免死循环。
                        logger.info("连续 %d 次未遵循协议，按最终回复处理", no_protocol_retries)
                        final_reply = response
                        break
                else:
                    # 没收到任何回复（等待超时 / 回复选择器未命中 / 站点结构变化）
                    self._emit(EventType.ERROR, {
                        "text": "未收到模型回复（可能是回复选择器未命中或站点结构变化）"
                    })
                    final_reply = "（未收到回复）"
                    break

            # 解析到协议命令，重置纠正计数
            no_protocol_retries = 0

            # 执行第一个指令
            cmd = cmds[0]
            self._emit(EventType.COMMAND_DETECTED, {"tool": cmd.command.tool, "id": cmd.id})
            self._emit(EventType.TOOL_START, {"tool": cmd.command.tool})

            # ── 任务范围守卫：拦截 AI 自行发明的用户未要求的新任务 ──
            guard_msg = self._scope_guard(cmd.command.tool, cmd.command.params, original_task)
            if guard_msg:
                logger.warning(f"[任务范围守卫] 拦截到越界工具调用: {cmd.command.tool} | {guard_msg[:120]}")
                current_input = guard_msg
                continue

            result = await self._execute_command(cmd.command)

            self._emit(EventType.TOOL_END, {
                "tool": cmd.command.tool,
                "status": result.status,
                "output": (result.output or result.error or "")[:6000],
            })
            self._emit(EventType.COMMAND_SUCCESS if result.status == "success" else EventType.COMMAND_ERROR, result)

            # 记录「动作-结果」摘要：回灌上下文的是这个，绝不回灌 agent 原话
            log_line = f"调用 {result.tool} → {result.status}: {(result.output or result.error)[:240]}"
            action_log.append(log_line)
            # 仅 DeepSeekSession 实现了 set_action_log；多平台会话（PlatformSession）没有，
            # 必须防护，否则切到通义/豆包等平台时会抛 AttributeError 导致整条命令崩溃。
            _sa = getattr(self._session, "set_action_log", None)
            if _sa is not None:
                try:
                    _sa(action_log)
                except Exception:
                    pass
            try:
                if self._memory:
                    self._memory.save(log_line, tags=["agent_action", f"turn_{turn}"])
            except Exception:
                pass

            # done() = AI 自己声明任务完成 → 循环正常结束。
            # 是否完成由 AI 自行判断，脚本不替它决定、也不强制 done()。
            if cmd.command.tool == "done":
                final_reply = self._strip_protocol(response) or last_ai_text or "[完成]"
                break

            # 构建下一轮输入：只回传工具的真实执行结果（成功/失败都如实），
            # 让 AI 自己判断是否已做完、是否该调用 done()。脚本不强制。
            current_input = (
                f"[系统] 工具 {result.tool} 执行结果:\n{result.output or result.error}\n\n"
                f"请继续；任务确实已完成时，调用 done() 结束。"
            )

        return final_reply or last_ai_text or "[完成]"

    async def run_with_loop(self, user_instruction: str, file_path: Optional[str] = None,
                            context_hints: str = "") -> str:
        """带自动循环的 run（和 run 相同，但名称更清晰）"""
        return await self.run(user_instruction, file_path, context_hints)

    async def continue_dialog(self) -> str:
        """
        继续对话（用户发送 '继续' 时调用）
        返回 AI 的回复内容
        """
        if self._correction_pending:
            correction = self._correction_pending
            self._correction_pending = None
            return await self.run(f"[系统纠正]\n{correction}")

        # 发送继续指令
        return await self.run("[系统] 用户要求继续，请接着上一步执行。")

    async def remember(self, content: str, tags: Optional[list] = None):
        """显式保存记忆"""
        if self._memory:
            mid = self._memory.save(content, tags=tags or [])
            return mid
        return None

    async def recall(self, query: str) -> list:
        """显式回忆记忆"""
        if self._memory:
            return self._memory.search(query)
        return []

    def set_memory_manager(self, mm: MemoryManager):
        """设置记忆管理器（外部注入）"""
        self._memory = mm

    def set_subagent_manager(self, sm):
        """设置子代理管理器（外部注入）"""
        self._subagent_manager = sm

    def stop(self):
        """停止运行循环"""
        self._running = False

    def inject_correction(self, correction_text: str):
        """注入系统纠正（AI 幻觉时人工干预）"""
        self._correction_pending = correction_text

    # ============================================================
    # 内部方法
    # ============================================================

    # 任务范围守卫：拦截「用户未明确要求的产物/调研/文件」类工具调用
    _DELIVERABLE_TOOLS = {
        "docx_create", "pptx_create",
        "browser_research", "browser_visit",
        "browser_search", "web_search", "summarize",
    }

    def _scope_guard(self, tool: str, params: dict, original_task: str) -> str:
        """若 AI 试图执行一个与用户原始指令无关的新任务（报告/调研/文件等），
        返回一段强制确认的系统指令；否则返回 None（放行）。"""
        if tool not in self._DELIVERABLE_TOOLS:
            return None
        arg = " ".join(str(v) for v in (params or {}).values() if isinstance(v, str))
        if not arg.strip() or not (original_task or "").strip():
            return None

        def _grams(s: str) -> set:
            s = re.sub(r"\s+", "", s or "")
            return set(s[i:i + 2] for i in range(len(s) - 1))

        og, ag = _grams(original_task), _grams(arg)
        if og and ag and (og & ag):
            return None  # 与原始任务有重叠，视为同一范畴，放行
        # 明显无关 → 禁止执行（但不许把任务踢回给用户，直接聚焦用户真正的要求）
        summary = arg[:80].replace("\n", " ")
        return (
            f"[SYSTEM] 你正准备执行一个用户【未明确要求】的任务（工具 {tool}，"
            f"参数摘要：「{summary}」），这与用户原始指令「{(original_task or '')[:120]}」无关。"
            f"禁止自行执行，也【不要向用户提问】。请忽略这件事，只专注于完成用户真正要求的任务；"
            f"你已把用户的任务做完时，调用 done() 结束即可。"
        )

    async def _execute_command(self, cmd: SimpleNamespace) -> ExecutionResult:
        """执行单个命令"""
        tool_name = getattr(cmd, "tool", None) or getattr(cmd, "action", None)
        params = getattr(cmd, "params", {}) or {}
        cmd_id = getattr(cmd, "id", str(uuid.uuid4()))

        if not tool_name:
            return ExecutionResult(id=cmd_id, status="error", error="指令缺少 tool 字段")

        # 特殊指令处理
        if tool_name == "done":
            return ExecutionResult(id=cmd_id, status="success", output="任务完成", tool="done")

        # ask / question 已被禁用：agent 没有「向用户提问」的权限，禁止偷懒把任务踢回给用户。
        # 若模型仍发出这类调用，明确拒绝并强制它用真实工具把任务推进到底。
        if tool_name in ("ask", "question"):
            return ExecutionResult(
                id=cmd_id, status="success",
                output=(
                    "[系统] 你【没有】 ask / question 工具，也不允许向用户提问或把任务踢回给用户。"
                    "遇到不确定或信息不全时，自己选择最合理的默认方案，直接用真实工具把用户的任务"
                    "做到底，完成后再调用 done()。现在请立即改用正确的工具继续推进。"
                ),
                tool=tool_name)

        # RAW 命令处理
        if tool_name == "raw_shell":
            return await self._tools.execute("shell_exec", {"command": params.get("command", "")})

        # 标准工具调用
        return await self._tools.execute(tool_name, params)

    def _emit(self, event_type: EventType, data=None):
        """触发事件回调"""
        if self._on_event:
            try:
                self._on_event(AgentEvent(event_type=event_type, data=data))
            except Exception as e:
                logger.warning(f"事件回调错误: {e}")

    @staticmethod
    def _strip_protocol(text: str) -> str:
        """去掉 @@@@ / <<<RAW>>> 协议块，保留自然语言部分作为最终回复。"""
        import re as _re
        t = _re.sub(r"@@@@.*?@@@@", "", text or "", flags=_re.DOTALL)
        t = _re.sub(r"<<<RAW>>>.*?<<<RAW>>>", "", t, flags=_re.DOTALL)
        return t.strip()

    def _safe_path(self, p) -> "Path":
        """解析用户指定的文件/目录路径：
        - 相对路径 → 拼到工作目录（D 盘）
        - **绝对路径：原样尊重用户指定**（用户叫在哪生成就在哪生成，
          包括 C 盘或其它盘，脚本绝不擅自重定向/改写路径）
        注意：agent 自身内部数据（profile/cookies/历史/任务/缓冲/子代理输出/
        Playwright 二进制）由 agent_core/xrz_paths.py 固定落在 D 盘，
        与「用户要求的产物输出路径」是两回事，互不影响。
        """
        from pathlib import Path as _P
        p = _P(p)
        if not p.is_absolute():
            return _P(self._work_dir) / p
        return p

    # ============================================================
    # 错误纠正（commander fix）
    # ============================================================

    async def fix(self, error_hint: str):
        """
        当检测到 AI 幻觉/错误时，注入系统级纠正
        使用方式：在 terminal.py 检测到异常后调用 commander.fix("纠正内容")
        """
        self.inject_correction(error_hint)
        return await self.continue_dialog()


# ============================================================
# 快捷函数（兼容旧代码）
# ============================================================

async def run_agent(browser_manager, session: DeepSeekSession, instruction: str, work_dir: str = "") -> str:
    """快捷函数：创建 Commander 并执行"""
    commander = Commander(
        browser_manager=browser_manager,
        session=session,
        work_dir=work_dir or os.getcwd(),
    )
    await commander.start(session=session)
    return await commander.run(instruction)
